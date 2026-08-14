#!/usr/bin/env python3
"""Build the AI Trust Platform / Platform Mesh deployment deck as a .pptx.
Self-contained (python-pptx only). Diagram slide left as a placeholder for the
externally-generated architecture image (Titan model retired)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

NAVY = RGBColor(0x0B, 0x2E, 0x4F)
TEAL = RGBColor(0x0E, 0x7C, 0x86)
LIGHT = RGBColor(0xF2, 0xF6, 0xF9)
GREY = RGBColor(0x44, 0x4B, 0x52)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x1F, 0x6F, 0xB2)
GREEN = RGBColor(0x2E, 0x8B, 0x57)

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _tf(box, size, color=GREY, bold=False, align=PP_ALIGN.LEFT):
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = ""
    return tf


def bg(slide, color=WHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def band(slide):
    bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(0.28))
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()


def title(slide, text, sub=None):
    band(slide)
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), SW - Inches(1.2), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = text
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = NAVY
    if sub:
        p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(15); r2.font.color.rgb = TEAL


def bullets(slide, items, top=1.7, left=0.7, width=None, size=16, gap=6):
    width = width or (SW - Inches(1.4))
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), width, SH - Inches(top) - Inches(0.5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        lvl = 0; txt = it; bold = False
        if isinstance(it, tuple):
            txt, lvl, bold = (it + (0, False))[:3]
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl; p.space_after = Pt(gap)
        r = p.add_run(); r.text = ("• " if lvl == 0 else "– ") + txt
        r.font.size = Pt(size - 2 * lvl); r.font.color.rgb = (NAVY if bold else GREY); r.font.bold = bold


def table(slide, rows, top=1.7, left=0.6, width=None, col_w=None, header=True, fs=13):
    width = width or (SW - Inches(1.2))
    nr = len(rows); nc = len(rows[0])
    h = Inches(0.42 * nr + 0.2)
    gt = slide.shapes.add_table(nr, nc, Inches(left), Inches(top), width, h).table
    if col_w:
        for j, w in enumerate(col_w): gt.columns[j].width = Inches(w)
    for i, row in enumerate(rows):
        for j, cell in enumerate(row):
            c = gt.cell(i, j); c.margin_left = Inches(0.08); c.margin_top = Inches(0.03); c.margin_bottom = Inches(0.03)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = c.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
            r = p.add_run(); r.text = str(cell)
            r.font.size = Pt(fs)
            if header and i == 0:
                r.font.bold = True; r.font.color.rgb = WHITE
                c.fill.solid(); c.fill.fore_color.rgb = NAVY
            else:
                r.font.color.rgb = GREY
                c.fill.solid(); c.fill.fore_color.rgb = (LIGHT if i % 2 else WHITE)
    return gt


def codebox(slide, text, top=1.7, left=0.6, width=None, height=None, fs=11):
    width = width or (SW - Inches(1.2)); height = height or (SH - Inches(top) - Inches(0.5))
    sp = slide.shapes.add_shape(1, Inches(left), Inches(top), width, height)
    sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor(0x0E, 0x1B, 0x2A); sp.line.fill.background()
    tf = sp.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.1)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line
        r.font.size = Pt(fs); r.font.name = "Consolas"; r.font.color.rgb = RGBColor(0xCF, 0xE8, 0xEF)


# ---- Slide 1: title ----
s = prs.slides.add_slide(BLANK); bg(s, NAVY)
tb = s.shapes.add_textbox(Inches(0.8), Inches(2.3), SW - Inches(1.6), Inches(2.6)); tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; r = p.add_run(); r.text = "Deploying the AI Trust Platform on Platform Mesh"
r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = WHITE
p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = "A high-level deployment workflow"
r2.font.size = Pt(22); r2.font.color.rgb = RGBColor(0x8F, 0xD3, 0xDB)
p3 = tf.add_paragraph(); p3.space_before = Pt(18); r3 = p3.add_run(); r3.text = "For PMs, Management & Developers"
r3.font.size = Pt(16); r3.font.color.rgb = RGBColor(0xB9, 0xC7, 0xD4)

# ---- Slide 2: one-liner ----
s = prs.slides.add_slide(BLANK); bg(s); title(s, "The one-liner")
card = s.shapes.add_shape(1, Inches(0.7), Inches(1.7), SW - Inches(1.4), Inches(1.6))
card.fill.solid(); card.fill.fore_color.rgb = LIGHT; card.line.color.rgb = TEAL
tf = card.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.3); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
r = tf.paragraphs[0].add_run(); r.text = ("We publish the AI Trust Platform once as a shared service on Platform Mesh. "
    "Any customer organization can then “subscribe” with one click and get their own fully "
    "isolated tenant — no new deployment per customer.")
r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = NAVY
bullets(s, [
    ("Management: faster onboarding, one thing to operate, strong per-customer data isolation for EU AI Act compliance.", 0, False),
    ("Developers: one shared app + a small operator that provisions a tenant on subscribe; isolation via per-org Keycloak realm + Postgres RLS + tenant-scoped ClickHouse/MinIO.", 0, False),
], top=3.6, size=16)

# ---- Slide 3: why this model ----
s = prs.slides.add_slide(BLANK); bg(s); title(s, "Why shared multi-tenant", "The core architectural decision (ADR-001)")
table(s, [
    ["Old model — copy per customer", "New model — shared multi-tenant"],
    ["Stamp ~23 services per customer", "ONE shared app for everyone"],
    ["Slow onboarding, high run cost", "Subscribe = seconds; one thing to operate"],
    ["Isolation = separate deployments", "Isolation = per-org identity + per-tenant data (RLS)"],
], top=1.8, col_w=[6.05, 6.05], fs=15)
bullets(s, [("Decision: the tenant = the Platform Mesh account. One account ⇄ one identity realm ⇄ one tenant_id — aligns 1:1 with the ApeiroRA / Platform Mesh account model.", 0, True)],
        top=4.4, size=15)

# ---- Slide 4: big picture (image, fallback to text diagram) ----
s = prs.slides.add_slide(BLANK); bg(s); title(s, "The big picture", "How two organizations stay isolated inside one app")
flow_img = "C:\\claude\\projects\\eu-ai-trust-platform\\aitrust_request_flow.png"
if os.path.exists(flow_img):
    s.shapes.add_picture(flow_img, Inches(1.35), Inches(1.7), width=Inches(10.6))
    bullets(s, [("Fail-closed: if the tenant can’t be verified, the request is rejected — never leaked.", 0, True)], top=6.55, size=13)
else:
    codebox(s, (
        "  Org A users                              Org B users\n"
        "      |  login (Keycloak realm A)              |  login (realm B)\n"
        "      |  verified token tenant_id=A            |  tenant_id=B\n"
        "      v                                        v\n"
        "      +----------->   ONE SHARED AI TRUST PLATFORM   <-----------+\n"
        "                        (runs once on Platform Mesh)\n"
        "                                 |\n"
        "                 tenant_id flows through every call\n"
        "        +------------------------+------------------------+\n"
        "        v                        v                        v\n"
        "   PostgreSQL (RLS:         ClickHouse               MinIO\n"
        "   own rows only)           (per-tenant traces)      (per-tenant files)\n"
    ), top=1.75, height=Inches(3.3), fs=12)
    bullets(s, [("Fail-closed: if the tenant can’t be verified, the request is rejected — never leaked.", 0, True)], top=5.3, size=15)

# ---- Slide 5: workflow phases ----
s = prs.slides.add_slide(BLANK); bg(s); title(s, "The deployment workflow", "10 steps, one command (deploy.sh) — grouped in 4 phases")
table(s, [
    ["Phase", "Steps", "What you get"],
    ["1. Prepare (once)", "0 check · 1 worker pool · 2 operator img · 2b app imgs", "Cluster ready + built images"],
    ["2. Publish (once)", "3 provider · 3b shared app", "Product in the marketplace + one shared app live"],
    ["3. Make consumable (once)", "4 consumer workspace · 5 bind APIs", "A customer org can see & subscribe"],
    ["4. Onboard (per customer)", "6 create subscription · 7 verify", "One-click tenant, verified live"],
], top=1.8, col_w=[3.2, 6.1, 3.0], fs=13)

# ---- Slide 6: what happens on Enable ----
s = prs.slides.add_slide(BLANK); bg(s); title(s, "What happens on “Enable”", "The subscription operator, automatically, in seconds")
bullets(s, [
    ("Reads the org from the subscription → that becomes the tenant.", 0, False),
    ("Creates a login client in the org’s Keycloak realm (stamps tenant_id on every token).", 0, False),
    ("Stands up a per-org front door (login proxy + URL ai-trust-mt-<org>…) into the shared app.", 0, False),
    ("Grants the org’s admin their roles.", 0, False),
    ("Reports the tenant’s URL back — ready to use.", 0, False),
    ("No new app is deployed — the tenant simply “turns on” inside the shared instance, data tagged & isolated from that moment.", 0, True),
], top=1.8, size=17, gap=10)

# ---- Slide 7: isolation guarantees ----
s = prs.slides.add_slide(BLANK); bg(s); title(s, "How isolation is guaranteed", "The trust story — outcome of a strict EU-auditor review (issue #16)")
table(s, [
    ["Layer", "Control", "Plain meaning"],
    ["Identity", "One Keycloak realm per org", "Org A users can’t log into Org B"],
    ["Token", "JWT cryptographically verified", "Nobody can forge “I’m tenant B”"],
    ["Database", "Postgres row-level security", "A query can’t return another tenant’s rows"],
    ["Telemetry", "ClickHouse filtered by tenant_id", "Traces/alerts are per-tenant"],
    ["Files", "MinIO keys prefixed per tenant", "Evidence files are per-tenant"],
    ["Everywhere", "Fail-closed", "Unknown tenant → deny, never leak"],
], top=1.75, col_w=[2.3, 4.5, 5.3], fs=13)

# ---- Slide 8: architecture diagram placeholder ----
s = prs.slides.add_slide(BLANK); bg(s); title(s, "Architecture diagram")
img = "C:\\claude\\projects\\eu-ai-trust-platform\\aitrust_mesh_architecture.png"
if os.path.exists(img):
    s.shapes.add_picture(img, Inches(1.2), Inches(1.7), width=Inches(11))
else:
    ph = s.shapes.add_shape(1, Inches(1.2), Inches(1.8), Inches(11), Inches(4.6))
    ph.fill.solid(); ph.fill.fore_color.rgb = LIGHT; ph.line.color.rgb = TEAL; ph.line.dash_style = 2
    tf = ph.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    r = tf.paragraphs[0].add_run(); r.text = "[ Insert architecture image: aitrust_mesh_architecture.png ]\n\n(Generate with the prompt in the notes, then drop it here.)"
    r.font.size = Pt(18); r.font.color.rgb = TEAL; tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# ---- Slide 9: roles ----
s = prs.slides.add_slide(BLANK); bg(s); title(s, "Roles & responsibilities")
table(s, [
    ["Who", "Owns"],
    ["Platform Mesh", "Account model, marketplace, per-org identity realms"],
    ["AI Trust Platform (this product)", "The shared app + subscription operator + data isolation"],
    ["The customer org", "Clicks Enable; manages its own users & data"],
], top=1.8, col_w=[4.0, 8.1], fs=15)
bullets(s, [("Scope boundary: the app enforces data + identity isolation; the mesh/deploy layer owns Kubernetes-level isolation & networking (documented in ADR-001).", 0, True)], top=4.2, size=15)

# ---- Slide 10: status ----
s = prs.slides.add_slide(BLANK); bg(s); title(s, "Status & what’s next")
bullets(s, [
    ("Done & live on staging (ai-trust-1):", 0, True),
    ("Shared app deployed; multi-tenancy enforced end-to-end; verified against a strict security audit.", 1, False),
    ("One-command deploy; one-click tenant onboarding; code pushed for PR review.", 1, False),
    ("Planned follow-ups (tracked, out of app scope):", 0, True),
    ("Kubernetes-resource-level isolation decision (namespace-per-tenant vs shared+RLS).", 1, False),
    ("Deeper Platform Mesh integration (declarative APIs, mesh-driven navigation).", 1, False),
    ("Operational hardening (credential rotation, CI/CD, SBOM/signing).", 1, False),
], top=1.8, size=16, gap=8)

# ---- Slide 11: developer deep-dive (1) — how it's built ----
s = prs.slides.add_slide(BLANK); bg(s); title(s, "For developers (1) — how it's built", "One shared app + a small operator; tenancy is a modular seam")
bullets(s, [
    ("libs/tenancy — the swappable seam (TENANCY_MODE = single | jwt | header; default single = no-op):", 0, True),
    ("resolver.py — in jwt mode, tenant comes ONLY from the VERIFIED token (PyJWT: RS256 sig via the realm's JWKS + exp + issuer must match TENANCY_JWKS_ISSUER_BASE). Client X-Tenant-Id is ignored.", 1, False),
    ("middleware.py — sets tenant_id_var (a ContextVar) per request; 401 fail-closed if unresolved.", 1, False),
    ("session.py — SQLAlchemy `begin` hook runs set_config('app.current_tenant', <tenant>) per transaction.", 1, False),
    ("Data-layer enforcement:", 0, True),
    ("Postgres RLS (migrations 0009/0010/0011): USING = own rows + catalog(NULL); WITH CHECK = write-own only; FORCE RLS. Runtime connects as non-superuser ai_trust_app (NOBYPASSRLS) — else RLS is inert.", 1, False),
    ("ClickHouse: tenant_clause() adds AND tenant_id={tenant} to every read (no RLS in CH); fail-closed 1=0 if no tenant. MinIO: object keys prefixed t/<tenant>/evidence/…", 1, False),
    ("The subscription operator (Go, controller-runtime) provisions a tenant on Enable — it does NOT deploy a new app.", 0, True),
], top=1.65, size=14, gap=6)

# ---- Slide 12: developer deep-dive (2) — request flow + operator + onboarding ----
s = prs.slides.add_slide(BLANK); bg(s); title(s, "For developers (2) — request flow & onboarding")
# left: request path codebox
codebox(s, (
    "REQUEST PATH (jwt mode)\n"
    "browser -> per-org oauth2-proxy\n"
    "   (auth vs org's Keycloak realm)\n"
    "   -> Authorization: Bearer <JWT>\n"
    "   -> nginx STRIPS client X-Tenant-Id\n"
    "backend middleware:\n"
    "   verify JWT (JWKS/exp/iss)\n"
    "   tenant_id_var.set(tenant)\n"
    "Postgres begin hook:\n"
    "   set_config('app.current_tenant')\n"
    "   -> RLS filters every row\n"
    "ClickHouse: AND tenant_id={t}\n"
    "MinIO: t/<tenant>/evidence/..."
), top=1.7, left=0.6, width=Inches(6.0), height=Inches(4.9), fs=11)
# right: operator + onboarding bullets
bullets(s, [
    ("Operator per Subscription reconcile:", 0, True),
    ("resolve org from spec.org (= tenant_id)", 1, False),
    ("create OIDC client in mesh realm <org> (+ tenant_id claim mapper)", 1, False),
    ("stamp per-org oauth2-proxy + Service + HTTPRoute (host ai-trust-mt-<org>…)", 1, False),
    ("seed org admin role tuple; write status.url", 1, False),
    ("finalizer soft-deletes proxy/route on unsubscribe (never data)", 1, False),
    ("Onboard a tenant:", 0, True),
    ("Subscription spec.org=<org> (+ adminEmail) via portal Enable or CR", 1, False),
    ("watch status → Ready, browse ai-trust-mt-<org>…  (no image build)", 1, False),
    ("Deploy = scripts 0→7; apply pg 0009→0011 + CH 0003", 1, False),
    ("Verified: 11 unit tests + RLS integration + live cross-tenant write denied", 1, False),
], top=1.7, left=7.0, width=Inches(5.7), size=13, gap=5)

out = "C:\\claude\\projects\\eu-ai-trust-platform\\AI_TRUST_PLATFORM_MESH_DEPLOYMENT.pptx"
try:
    prs.save(out)
except PermissionError:
    print("LOCKED: the primary .pptx is still open in PowerPoint — close it and re-run. NOT saved.")
    raise SystemExit(1)
print("SAVED", out, "slides:", len(prs.slides._sldIdLst))
